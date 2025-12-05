"""
Document parsing module.
Handles text extraction from various file formats.
"""

import os
from typing import Optional

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    try:
        import PyPDF2
        text_parts = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
                except Exception as e:
                    text_parts.append(f"--- Page {page_num + 1} ---\n[Error extracting page: {e}]")
        return "\n\n".join(text_parts) if text_parts else "[No text content extracted from PDF]"
    except ImportError:
        return "[Error: PyPDF2 not installed]"
    except Exception as e:
        return f"[Error extracting PDF: {str(e)}]"

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs) if paragraphs else "[No text content extracted from DOCX]"
    except ImportError:
        return "[Error: python-docx not installed]"
    except Exception as e:
        return f"[Error extracting DOCX: {str(e)}]"

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts) if text_parts else "[No text content extracted from PPTX]"
    except ImportError:
        return "[Error: python-pptx not installed]"
    except Exception as e:
        return f"[Error extracting PPTX: {str(e)}]"

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text or markdown file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            return f"[Error reading text file: {str(e)}]"
    except Exception as e:
        return f"[Error reading text file: {str(e)}]"

def extract_text(file_path: str, extension: str) -> str:
    """Route text extraction to appropriate handler."""
    extension = extension.lower()
    
    if extension in ['.png', '.jpg', '.jpeg', '.gif', '.webp']:
        file_size = os.path.getsize(file_path)
        return f"[Image file: {os.path.basename(file_path)} - {file_size} bytes. No OCR performed.]"
        
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension == '.docx':
        return extract_text_from_docx(file_path)
    elif extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif extension in ['.txt', '.md']:
        return extract_text_from_txt(file_path)
    else:
        return f"[Unsupported file type: {extension}]"
